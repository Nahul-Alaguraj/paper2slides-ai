from pptx import Presentation
from pptx.util import Inches, Pt
from typing import List, Dict
import uuid
import os


def generate_ppt(slides: List[Dict], title: str = "AI Generated Slides") -> str:
    prs = Presentation()

    # ---- Title slide ----
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated using AI Paper-to-Slides Generator"

    # ---- Content slides ----
    content_layout = prs.slide_layouts[1]

    for s in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = s.get("title", "Slide")

        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()

        bullets = s.get("bullets", [])

        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)

    # ---- Save file ----
    os.makedirs("generated_ppts", exist_ok=True)
    filename = f"slides_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join("generated_ppts", filename)

    prs.save(filepath)
    return filepath
